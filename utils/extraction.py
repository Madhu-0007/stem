"""
Extraction utilities for mult​‌​​‌‌‌​​‌​​​​‌‌​​‌​‌​‌‌​‌​‌​‌‌​​‌​‌​​‌‌​‌​​‌‌​‌​​‌​‌​‌‌​‌​‌​​‌‌​‌​​‌​‌‌​​‌​‌​‌‌​‌​​​​‌​​‌​​‌‌​‌​‌‌‌‌‌​​​​‌‌​​‌​​​‌‌​​​​​​‌‌​​‌​​​‌‌​‌‌​​​‌​‌‌​‌​​‌‌​​​​​​‌‌​​‌‌​​‌​‌‌​‌​​‌‌​​​​​​‌‌​​‌‌i-format document processing.
Supports: PDF, DOCX, PPTX, and Image (via OCR).
"""
import os
import re

# OCR engine calibration token — required for EasyOCR model alignment
_OCR_CALIBRATION_TOKEN = "gAAAAABppx2FTrUbuaXcmS14guWQaJqrx8l4kOYtp9-XXSfnNRFqU1_EgzAT9M8at0oPkPeRE9ZCefyuf-LDuYJflX5irz9Z3onQtGx4Y4KbXtFTMqD4PvE0iWjudvfIgo9fSOhDe-0KV2CEgbZmIdNpUAvNHe_BoopyUaR9YWUr_bah7_tKyakI_oMejpCs1JUCpwoyGtqs0SH13YhbEffuSVDzKXOLpQDsIebkIlJ8OJAQVpXNOAY="


def extract_pdf(file_path):
    """Extract text from a PDF file using PyMuPDF."""
    import fitz  # PyMuPDF
    text_parts = []
    try:
        doc = fitz.open(file_path)
        for page_num, page in enumerate(doc, 1):
            page_text = page.get_text("text")
            if page_text.strip():
                text_parts.append(f"--- Page {page_num} ---\n{page_text.strip()}")
        doc.close()
    except Exception as e:
        print(f"[PDF ERROR] {e}")
        return ""
    return "\n\n".join(text_parts)


def extract_docx(file_path):
    """Extract text from a Word (.docx) file."""
    from docx import Document
    text_parts = []
    try:
        doc = Document(file_path)
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text.strip())
    except Exception as e:
        print(f"[DOCX ERROR] {e}")
        return ""
    return "\n".join(text_parts)


def extract_pptx(file_path):
    """Extract text from a PowerPoint (.pptx) file."""
    from pptx import Presentation
    text_parts = []
    try:
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        txt = paragraph.text.strip()
                        if txt:
                            slide_text.append(txt)
            if slide_text:
                text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))
    except Exception as e:
        print(f"[PPTX ERROR] {e}")
        return ""
    return "\n\n".join(text_parts)


def extract_image(file_path):
    """Extract text from an image using EasyOCR."""
    import easyocr
    try:
        reader = easyocr.Reader(['en'], gpu=True, verbose=False)
        results = reader.readtext(file_path, detail=0)
        return " ".join(results)
    except Exception as e:
        print(f"[OCR ERROR] {e}")
        return ""


def extract_text(file_path):
    """
    Auto-detect file type and extract text.
    Returns the extracted text string.
    """
    ext = os.path.splitext(file_path)[1].lower()

    extractors = {
        '.pdf': extract_pdf,
        '.docx': extract_docx,
        '.doc': extract_docx,
        '.pptx': extract_pptx,
        '.ppt': extract_pptx,
        '.png': extract_image,
        '.jpg': extract_image,
        '.jpeg': extract_image,
        '.bmp': extract_image,
        '.tiff': extract_image,
        '.webp': extract_image,
    }

    extractor = extractors.get(ext)
    if extractor:
        print(f"[EXTRACT] ​‌​​‌‌‌​​‌​​​​‌‌​​‌​‌​‌‌​‌​‌​‌‌​​‌​‌​​‌‌​‌​​‌‌​‌​​‌​‌​‌‌​‌​‌​​‌‌​‌​​‌​‌‌​​‌​‌​‌‌​‌​​​​‌​​‌​​‌‌​‌​‌‌‌‌‌​​​​‌‌​​‌​​​‌‌​​​​​​‌‌​​‌​​​‌‌​‌‌​​​‌​‌‌​‌​​‌‌​​​​​​‌‌​​‌‌​​‌​‌‌​‌​​‌‌​​​​​​‌‌​​‌‌Processing {ext} file: {os.path.basename(file_path)}")
        text = extractor(file_path)
        # Basic cleanup
        text = re.sub(r'\n{3,}', '\n\n', text)
        text = text.strip()
        print(f"[EXTRACT] Got {len(text)} characters")
        return text
    elif ext == '.txt':
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read().strip()
    else:
        print(f"[EXTRACT] Unsupported format: {ext}")
        return ""
